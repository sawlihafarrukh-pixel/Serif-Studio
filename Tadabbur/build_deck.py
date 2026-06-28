#!/usr/bin/env python3
"""
Build an editable PowerPoint deck on Tadabbur, based on the chapter from
Ayesha Syahira's book "Befriending the Quran".

All text lives in real editable text boxes / placeholders. No images embedded
so it stays clean and fully editable. Run: python3 build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ARABIC_FONT = "Geeza Pro"  # macOS naskh; legible with tashkeel

# ---------------------------------------------------------------- palette ----
# Calm, study-friendly palette (independent of Serif Studio brand).
INK        = RGBColor(0x2B, 0x2B, 0x2B)   # near-black body text
DEEP       = RGBColor(0x1F, 0x3A, 0x34)   # deep green — headers
ACCENT     = RGBColor(0xB9, 0x7A, 0x3E)   # warm gold — accents/numbers
MUTED      = RGBColor(0x6B, 0x6B, 0x6B)   # captions/footnotes
CREAM      = RGBColor(0xF6, 0xF3, 0xEC)   # background
CARD       = RGBColor(0xFF, 0xFF, 0xFF)   # card fill
SOFT       = RGBColor(0xE7, 0xDF, 0xD2)   # soft divider/card border

# 16:9
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM
    return s


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def para(tf, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT,
         space_after=8, space_before=0, level=0, first=False, font="Calibri"):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.add_run()
    set_run(r, text, size, color, bold, italic, font)
    return p


def arabic_para(tf, text, size, color, align=PP_ALIGN.RIGHT, bold=False,
                space_after=8, space_before=0, first=False):
    """Right-to-left paragraph with an Arabic complex-script font set."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p._p.get_or_add_pPr().set('rtl', '1')   # mark paragraph RTL
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = ARABIC_FONT               # latin typeface
    rPr = r._r.get_or_add_rPr()
    cs = rPr.find(qn('a:cs'))               # complex-script typeface (Arabic)
    if cs is None:
        cs = rPr.makeelement(qn('a:cs'), {})
        rPr.append(cs)
    cs.set('typeface', ARABIC_FONT)
    return p


def accent_bar(slide, left, top, width=Inches(1.1), height=Pt(5), color=ACCENT):
    from pptx.enum.shapes import MSO_SHAPE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def card(slide, left, top, width, height, fill=CARD, border=SOFT):
    from pptx.enum.shapes import MSO_SHAPE
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = border
    c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


_page = {"n": 2}  # slide 1 (title) carries no footer

def footer(slide, page=None):
    p = page if page is not None else _page["n"]
    _page["n"] = p + 1
    tb, tf = textbox(slide, Inches(0.5), Inches(7.02), Inches(9), Inches(0.4))
    para(tf, "Tadabbur — reflecting on the Qur’an  ·  after Ayesha Syahira, “Befriending the Qur’an”",
         9, MUTED, first=True)
    tb2, tf2 = textbox(slide, Inches(12.2), Inches(7.02), Inches(0.8), Inches(0.4))
    para(tf2, str(p), 9, MUTED, align=PP_ALIGN.RIGHT, first=True)


# ============================================================ SLIDE 1: TITLE
s = add_slide()
# deep band
from pptx.enum.shapes import MSO_SHAPE
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.35), SW, Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb = DEEP
band.line.fill.background(); band.shadow.inherit = False

tb, tf = textbox(s, Inches(1.0), Inches(2.55), Inches(11.3), Inches(1.4))
para(tf, "Tadabbur", 60, CREAM, bold=True, first=True)
para(tf, "reflecting on the Qur’an — and letting it change you", 24, RGBColor(0xD9,0xCB,0xB0), italic=True, space_before=4)

tb, tf = textbox(s, Inches(1.0), Inches(5.45), Inches(11.3), Inches(1.0))
para(tf, "A short talk based on the chapter from", 14, MUTED, first=True)
para(tf, "“Befriending the Qur’an” — Ayesha Syahira", 16, DEEP, bold=True, space_before=2)
accent_bar(s, Inches(1.02), Inches(2.4), width=Inches(1.4), color=ACCENT)

# ============================================================ SLIDE 2: HOOK
s = add_slide()
tb, tf = textbox(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(2.8))
para(tf, "You finish a whole juz’…", 40, DEEP, bold=True, first=True, space_after=10)
para(tf, "…and a minute later, you couldn’t say what it was about.", 30, INK, italic=True, space_after=24)
para(tf, "Sound familiar?", 22, ACCENT, bold=True)
accent_bar(s, Inches(1.02), Inches(2.05), width=Inches(1.4), color=ACCENT)
tb, tf = textbox(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.8))
para(tf, "Reciting is not the same as understanding. That gap is what tadabbur is for.",
     18, MUTED, first=True)
footer(s)

# ====================================================== SLIDE 3: WHY REFLECT
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "Why reflect at all?", 34, DEEP, bold=True, first=True)

# Arabic aayah — Muhammad 47:24
tb, tf = textbox(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.3))
arabic_para(tf, "أَفَلَا يَتَدَبَّرُونَ الْقُرْآنَ أَمْ عَلَىٰ قُلُوبٍ أَقْفَالُهَا",
            40, DEEP, first=True)
# English + citation
tb, tf = textbox(s, Inches(0.9), Inches(3.5), Inches(11.5), Inches(1.4))
para(tf, "“Then do they not reflect upon the Qur’an,\nor are there locks upon [their] hearts?”",
     26, INK, italic=True, first=True, space_after=8)
para(tf, "— Muhammad 47:24", 18, ACCENT, bold=True)

tb, tf = textbox(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.1))
para(tf, "We aren’t asked only to recite the Qur’an — we’re asked to ponder it. "
         "Reflection is what turns reading into a relationship.", 18, MUTED, first=True)
footer(s)

# ================================================ SLIDE 4: AN-NISA 4:82 (verse)
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "And it holds together", 34, DEEP, bold=True, first=True)

# Arabic aayah — an-Nisa 4:82 (longer verse)
tb, tf = textbox(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(1.7))
arabic_para(tf, "أَفَلَا يَتَدَبَّرُونَ الْقُرْآنَ ۚ وَلَوْ كَانَ مِنْ عِندِ "
                "غَيْرِ اللَّهِ لَوَجَدُوا فِيهِ اخْتِلَافًا كَثِيرًا",
            34, DEEP, first=True)
tb, tf = textbox(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(1.6))
para(tf, "“Then do they not reflect upon the Qur’an? If it had been from [any] other "
         "than Allah, they would have found within it much contradiction.”",
     22, INK, italic=True, first=True, space_after=8)
para(tf, "— an-Nisa 4:82", 18, ACCENT, bold=True)

tb, tf = textbox(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.9))
para(tf, "Reflection doesn’t just deepen faith — it’s where the Qur’an proves itself.",
     18, MUTED, first=True)
footer(s)

# =================================================== SLIDE 3: WHAT IS TADABBUR
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "What is tadabbur?", 34, DEEP, bold=True, first=True)

card(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.7))
tb, tf = textbox(s, Inches(1.25), Inches(2.25), Inches(10.8), Inches(1.3))
para(tf, "Tadabbur means pondering and reflecting upon the Qur’an — ", 22, INK, first=True)
p = tf.paragraphs[0]
r = p.add_run(); set_run(r, "thinking deeply about the meaning of the aayaat and how they relate to your own life.", 22, INK, italic=True)

tb, tf = textbox(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(2.6))
para(tf, "Two things to hold onto:", 18, DEEP, bold=True, first=True, space_after=12)
para(tf, "•  the Qur’an’s own word — yatadabbaroon, “do they not reflect?” — is a question. An invitation.", 19, INK, space_after=12)
para(tf, "•  the goal isn’t just to know the aayah — it’s to be moved by it.", 19, INK)
footer(s)

# ============================================ SLIDE 4: 5 PRINCIPLES (overview)
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "Five principles for seeking knowledge", 32, DEEP, bold=True, first=True)
tb, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
para(tf, "The author’s posture for seeking knowledge — the ground tadabbur grows from.", 16, MUTED, first=True)

items = [
    ("1", "Set sincere intentions", "tie every intention back to Allah"),
    ("2", "Keep an open mind", "empty the cup; stay humble and ready to learn"),
    ("3", "Have an action plan", "knowledge is meant to be lived, not stored"),
    ("4", "Revise & pass it on", "review your notes; share what you’ve learned"),
    ("5", "Don’t give up", "the journey is long — keep going, ask Allah"),
]
top = Inches(2.5)
row_h = Inches(0.82)
for i, (num, title, sub) in enumerate(items):
    y = Emu(int(top) + i * int(row_h))
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y, Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background(); circle.shadow.inherit = False
    ctf = circle.text_frame; ctf.word_wrap = True
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); set_run(cr, num, 22, RGBColor(0xFF,0xFF,0xFF), bold=True)
    tb, tf = textbox(s, Inches(1.8), Emu(int(y) - Emu(int(Inches(0.05)))), Inches(10.4), Inches(0.8))
    para(tf, title, 20, DEEP, bold=True, first=True, space_after=0)
    p = tf.add_paragraph(); rr = p.add_run(); set_run(rr, sub, 14, MUTED)
footer(s)

# ===================================== SLIDES 5-9: each principle, one per slide
principles = [
    ("1", "Set sincere intentions",
     "Before you seek knowledge, fix why you’re seeking it.",
     ["Tie your intention back to Allah — not to praise, status or habit.",
      "A sincere intention turns ordinary study into worship."]),
    ("2", "Keep an open mind",
     "Come to the Qur’an ready to be taught.",
     ["Humble yourself — empty the cup before you try to fill it.",
      "Be willing to trade an old understanding for a truer one."]),
    ("3", "Have an action plan",
     "Knowledge that isn’t acted on isn’t complete.",
     ["For every reflection, ask: what will I actually do?",
      "Action is what brings the learning to fruition."]),
    ("4", "Revise what you’ve learned — and pass it on",
     "Protect knowledge by revisiting and sharing it.",
     ["Notes fade from memory — revise so they don’t fade away.",
      "Teaching others completes the circle and deepens you."]),
    ("5", "Don’t give up",
     "Some days nothing lands. This is for those days.",
     ["Don’t let doubt make the knowledge melt away.",
      "Ask Allah for barakah; keep pulling yourself closer to Him."]),
]
for num, title, lead, bullets in principles:
    s = add_slide()
    # big ghost number
    tb, tf = textbox(s, Inches(9.7), Inches(0.6), Inches(3.3), Inches(3.0))
    para(tf, num, 150, SOFT, bold=True, align=PP_ALIGN.RIGHT, first=True)
    accent_bar(s, Inches(0.9), Inches(0.7))
    tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(9.0), Inches(0.55))
    para(tf, f"PRINCIPLE {num} OF 5", 13, ACCENT, bold=True, first=True)
    tb, tf = textbox(s, Inches(0.85), Inches(1.45), Inches(8.6), Inches(1.6))
    para(tf, title, 32, DEEP, bold=True, first=True)
    tb, tf = textbox(s, Inches(0.9), Inches(3.0), Inches(9.2), Inches(0.9))
    para(tf, lead, 20, ACCENT, italic=True, first=True)
    tb, tf = textbox(s, Inches(0.9), Inches(4.1), Inches(9.6), Inches(2.6))
    first = True
    for b in bullets:
        para(tf, "•  " + b, 19, INK, first=first, space_after=16)
        first = False
    footer(s)

# ============================================ SLIDE 10: THE METHOD (4 columns)
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "How to do your tadabbur", 34, DEEP, bold=True, first=True)
tb, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
para(tf, "A simple four-column note for each aayah you reflect on.", 16, MUTED, first=True)

cols = [
    ("1", "Write the aayah", "Write the aayah and its meaning. In time, learn it in Arabic too."),
    ("2", "Reflect & relate", "Who is it addressing? What does it teach? How does it relate to your life right now?"),
    ("3", "Action plan", "Turn the reflection into practice — one concrete thing you’ll do."),
    ("4", "Du’aa", "Close by asking Allah — a du’aa inspired by the aayah."),
]
cw = Inches(2.78)
gap = Inches(0.18)
left0 = Inches(0.9)
top = Inches(2.55)
ch = Inches(3.7)
for i, (num, head, body) in enumerate(cols):
    x = Emu(int(left0) + i * (int(cw) + int(gap)))
    card(s, x, top, cw, ch)
    # number chip
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(x)+int(Inches(0.25))), Emu(int(top)+int(Inches(0.25))), Inches(0.55), Inches(0.55))
    chip.fill.solid(); chip.fill.fore_color.rgb = DEEP
    chip.line.fill.background(); chip.shadow.inherit = False
    cp = chip.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); set_run(cr, num, 20, CREAM, bold=True)
    tb, tf = textbox(s, Emu(int(x)+int(Inches(0.22))), Emu(int(top)+int(Inches(1.05))), Emu(int(cw)-int(Inches(0.44))), Inches(2.5))
    para(tf, head, 18, DEEP, bold=True, first=True, space_after=8)
    para(tf, body, 14, INK)
tb, tf = textbox(s, Inches(0.9), Inches(6.42), Inches(11.5), Inches(0.55))
para(tf, "Try it now: take one aayah you read this week and run it through the four columns.",
     15, ACCENT, italic=True, bold=True, first=True)
footer(s)

# ================================ SLIDE: DECORATE / RETURN (light, merged point)
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "Make notes you’ll come back to", 34, DEEP, bold=True, first=True)
tb, tf = textbox(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.6))
para(tf, "Don’t just write — decorate.", 28, INK, bold=True, first=True, space_after=18)
para(tf, "Colour and care make notes you’ll actually reread — and the writing itself "
         "is therapeutic. Notes you enjoy are notes you return to.", 22, MUTED, italic=True)
footer(s)

# =================================================== SLIDE: BUILD A CADENCE
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.7))
tb, tf = textbox(s, Inches(0.85), Inches(0.85), Inches(11.6), Inches(1.0))
para(tf, "Turn it into a rhythm", 34, DEEP, bold=True, first=True)
tb, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
para(tf, "Set Qur’an goals across several phases — small habits that compound.", 16, MUTED, first=True)

cadence = [
    ("Daily", "A set number of pages, with a few minutes for tadabbur."),
    ("Weekly", "A class or lecture that explains what you’re reading."),
    ("Monthly", "A bigger goal — e.g. a surah to learn or memorise."),
    ("Yearly", "The big aim that all the smaller goals add up to."),
    ("Ramadhan", "Build on the habit and push for more in the blessed month."),
]
top = Inches(2.5)
row_h = Inches(0.86)
for i, (label, desc) in enumerate(cadence):
    y = Emu(int(top) + i*int(row_h))
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), y, Inches(2.4), Inches(0.66))
    chip.fill.solid(); chip.fill.fore_color.rgb = DEEP
    chip.line.fill.background(); chip.shadow.inherit = False
    cp = chip.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    chip.text_frame.word_wrap = True
    cr = cp.add_run(); set_run(cr, label, 18, CREAM, bold=True)
    tb, tf = textbox(s, Inches(3.6), Emu(int(y)+int(Inches(0.07))), Inches(9.0), Inches(0.7))
    para(tf, desc, 18, INK, first=True)
footer(s)

# ===================================================== SLIDE 13: ONE TAKEAWAY
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
band.fill.solid(); band.fill.fore_color.rgb = DEEP
band.line.fill.background(); band.shadow.inherit = False
accent_bar(s, Inches(1.0), Inches(2.0), width=Inches(1.4), color=ACCENT)
tb, tf = textbox(s, Inches(1.0), Inches(2.3), Inches(11.3), Inches(3.0))
para(tf, "Read less, reflect more.", 46, CREAM, bold=True, first=True, space_after=20)
para(tf, "One aayah. One reflection. One action. One du’aa.", 24, RGBColor(0xD9,0xCB,0xB0), italic=True)
tb, tf = textbox(s, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.6))
para(tf, "based on “Befriending the Qur’an” — Ayesha Syahira", 14, RGBColor(0xAF,0xA0,0x86), first=True)

# ===================================================== SLIDE 14: THANK YOU
s = add_slide()
tb, tf = textbox(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.6))
para(tf, "Shukran — thank you", 40, DEEP, bold=True, align=PP_ALIGN.CENTER, first=True)
para(tf, "May Allah make us among those who reflect, and act.", 20, MUTED, italic=True, align=PP_ALIGN.CENTER, space_before=8)
accent_bar(s, Inches(6.0), Inches(4.35), width=Inches(1.3), color=ACCENT)
tb, tf = textbox(s, Inches(1.0), Inches(4.7), Inches(11.3), Inches(1.4))
para(tf, "Start tonight: one aayah, the four columns.", 22, DEEP, bold=True, align=PP_ALIGN.CENTER, first=True, space_after=6)
para(tf, "Read more: “Befriending the Qur’an” — Ayesha Syahira", 15, MUTED, align=PP_ALIGN.CENTER)

out = "Tadabbur.pptx"
prs.save(out)
print("saved", out, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
